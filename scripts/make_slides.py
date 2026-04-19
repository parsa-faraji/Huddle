"""Generate a beginner-friendly walkthrough deck for Huddle.

Run: python3 scripts/make_slides.py
Output: huddle-walkthrough.pptx in the repo root.
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Huddle palette
AMBER_DEEP = RGBColor(0xFF, 0xB0, 0x00)
AMBER_MID = RGBColor(0xFF, 0xDC, 0x90)
AMBER_SOFT = RGBColor(0xFF, 0xEC, 0xC1)
NAVY = RGBColor(0x0C, 0x1A, 0x35)
CHARCOAL = RGBColor(0x1A, 0x1A, 0x2E)
BROWN = RGBColor(0x5C, 0x40, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x94, 0x88, 0x70)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    size=24,
    bold=False,
    color=CHARCOAL,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font="Helvetica Neue",
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return tb


def add_bullets(
    slide,
    left,
    top,
    width,
    height,
    bullets,
    *,
    size=18,
    color=CHARCOAL,
    bullet_color=AMBER_DEEP,
    font="Helvetica Neue",
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        p.alignment = PP_ALIGN.LEFT
        r_bullet = p.add_run()
        r_bullet.text = "●  "
        r_bullet.font.size = Pt(size)
        r_bullet.font.color.rgb = bullet_color
        r_bullet.font.bold = True
        r_bullet.font.name = font
        r_text = p.add_run()
        r_text.text = item
        r_text.font.size = Pt(size)
        r_text.font.color.rgb = color
        r_text.font.name = font
    return tb


def add_code(slide, left, top, width, height, code, *, size=14):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY
    box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    for i, line in enumerate(code.splitlines()):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = "Menlo"
        r.font.size = Pt(size)
        r.font.color.rgb = WHITE
    return box


def add_pill(slide, left, top, width, height, text, *, fill=AMBER_DEEP, text_color=CHARCOAL, size=14, bold=True):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.fill.background()
    box.shadow.inherit = False
    # Force rounded corners
    try:
        box.adjustments[0] = 0.5
    except Exception:
        pass
    tf = box.text_frame
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = text_color
    r.font.name = "Helvetica Neue"
    return box


def add_header_band(slide, step_label, title):
    # Amber band across the top
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.1))
    band.fill.solid()
    band.fill.fore_color.rgb = AMBER_MID
    band.line.fill.background()
    band.shadow.inherit = False
    if step_label:
        add_pill(slide, Inches(0.5), Inches(0.25), Inches(1.6), Inches(0.5), step_label, fill=NAVY, text_color=AMBER_MID, size=14)
        add_text(slide, Inches(2.2), Inches(0.2), Inches(10.5), Inches(0.7), title, size=30, bold=True, color=CHARCOAL)
    else:
        add_text(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.7), title, size=32, bold=True, color=CHARCOAL)


def add_footer(slide, index, total):
    add_text(
        slide,
        Inches(0.5),
        Inches(7.05),
        Inches(6),
        Inches(0.35),
        "Huddle — build a campus study app",
        size=10,
        color=BROWN,
    )
    add_text(
        slide,
        Inches(11.5),
        Inches(7.05),
        Inches(1.5),
        Inches(0.35),
        f"{index} / {total}",
        size=10,
        color=BROWN,
        align=PP_ALIGN.RIGHT,
    )


# -------------------- Slide builders --------------------


def title_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s, AMBER_SOFT)
    # Decorative circle
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(-2), Inches(8), Inches(8))
    c.fill.solid()
    c.fill.fore_color.rgb = AMBER_DEEP
    c.line.fill.background()
    c.shadow.inherit = False
    # Small navy dot
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(5), Inches(3), Inches(3))
    d.fill.solid()
    d.fill.fore_color.rgb = NAVY
    d.line.fill.background()
    d.shadow.inherit = False

    add_text(s, Inches(0.8), Inches(2.3), Inches(9), Inches(1.2), "Huddle", size=80, bold=True, color=CHARCOAL, font="Georgia")
    add_text(s, Inches(0.8), Inches(3.4), Inches(9), Inches(0.7), "Build a campus study-spot app from scratch", size=26, color=CHARCOAL)
    add_text(
        s,
        Inches(0.8),
        Inches(4.3),
        Inches(8.5),
        Inches(1),
        "React · Vite · Tailwind · Firebase · Leaflet",
        size=18,
        color=BROWN,
    )
    add_pill(s, Inches(0.8), Inches(5.4), Inches(3.2), Inches(0.55), "Beginner friendly", fill=NAVY, text_color=AMBER_MID, size=14)
    return s


def overview_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s, WHITE)
    add_header_band(s, "Overview", "What you'll build")
    add_text(
        s,
        Inches(0.6),
        Inches(1.4),
        Inches(7.5),
        Inches(0.5),
        "A phone-sized web app where students:",
        size=18,
        color=CHARCOAL,
    )
    add_bullets(
        s,
        Inches(0.6),
        Inches(2),
        Inches(7.8),
        Inches(4.5),
        [
            "Sign up with email + password",
            "Browse campus study spots in a list or on a map",
            "Filter by quiet / outlets / open now",
            "Rate spots across noise, seating, WiFi, lighting, crowdedness…",
            "Set personal preferences and see recommended spots",
            "Create and join study groups by course",
            "Get group-level spot recommendations based on member preferences",
        ],
        size=17,
    )
    # Right column: visual block
    v = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9), Inches(1.5), Inches(3.8), Inches(5))
    v.fill.solid()
    v.fill.fore_color.rgb = AMBER_MID
    v.line.fill.background()
    v.shadow.inherit = False
    add_text(s, Inches(9), Inches(1.7), Inches(3.8), Inches(0.6), "Huddle", size=28, bold=True, color=CHARCOAL, align=PP_ALIGN.CENTER, font="Georgia")
    add_text(s, Inches(9), Inches(2.3), Inches(3.8), Inches(0.45), "Find a study spot!", size=14, color=CHARCOAL, align=PP_ALIGN.CENTER)
    # mock cards
    for i, name in enumerate(["Doe Library", "MLK Student Union", "Cafe Strada"]):
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(2.9 + i * 1), Inches(3.4), Inches(0.9))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = AMBER_DEEP
        card.shadow.inherit = False
        add_text(s, Inches(9.4), Inches(3 + i * 1), Inches(3.2), Inches(0.3), name, size=13, bold=True, color=CHARCOAL)
        add_text(s, Inches(9.4), Inches(3.3 + i * 1), Inches(3.2), Inches(0.3), "0.3 mi · Silent · Open", size=10, color=BROWN)
    return s


def stack_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s, WHITE)
    add_header_band(s, "Tech stack", "The tools and what each one does")
    cells = [
        ("React", "UI library — components that update when data changes", AMBER_DEEP),
        ("Vite", "Fast dev server + production bundler", AMBER_MID),
        ("Tailwind", "Utility-first styling, no CSS files", AMBER_DEEP),
        ("Firebase Auth", "Email/password login with zero backend code", AMBER_MID),
        ("Firestore", "Real-time database — data updates appear instantly", AMBER_DEEP),
        ("Leaflet + OSM", "Free interactive maps, no API key needed", AMBER_MID),
        ("React Router", "Move between pages without full reloads", AMBER_DEEP),
        ("Vercel", "Deploy the whole thing with one command", AMBER_MID),
    ]
    for i, (name, desc, color) in enumerate(cells):
        col = i % 2
        row = i // 2
        left = Inches(0.6 + col * 6.1)
        top = Inches(1.4 + row * 1.4)
        pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.8), Inches(0.55))
        pill.fill.solid()
        pill.fill.fore_color.rgb = color
        pill.line.fill.background()
        pill.shadow.inherit = False
        try:
            pill.adjustments[0] = 0.5
        except Exception:
            pass
        tf = pill.text_frame
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = name
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = CHARCOAL
        r.font.name = "Helvetica Neue"
        add_text(s, left + Inches(2), top + Inches(0.1), Inches(4.0), Inches(0.5), desc, size=13, color=CHARCOAL)
    return s


def prereqs_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s, WHITE)
    add_header_band(s, "Step 0", "Tools you'll install once")
    add_bullets(
        s,
        Inches(0.6),
        Inches(1.4),
        Inches(12),
        Inches(3),
        [
            "Node.js 18 or newer (from nodejs.org)",
            "A code editor — VS Code is free and great",
            "A free Google account (for Firebase) and GitHub account (for Vercel)",
            "A terminal — Terminal on Mac, PowerShell on Windows",
        ],
        size=18,
    )
    add_text(s, Inches(0.6), Inches(4.4), Inches(12), Inches(0.4), "Quick check: run these in your terminal", size=14, color=BROWN, bold=True)
    add_code(
        s,
        Inches(0.6),
        Inches(4.8),
        Inches(12),
        Inches(1.6),
        "node -v      # should print v18.x or newer\nnpm -v       # should print 9.x or newer\ngit --version",
        size=15,
    )
    return s


def step_scaffold(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s, WHITE)
    add_header_band(s, "Step 1", "Scaffold the project with Vite")
    add_text(s, Inches(0.6), Inches(1.35), Inches(12), Inches(0.4), "Create the app in one command. Vite sets up React + build tooling.", size=15, color=BROWN)
    add_code(
        s,
        Inches(0.6),
        Inches(1.9),
        Inches(12),
        Inches(2),
        "npm create vite@latest huddle -- --template react\ncd huddle\nnpm install\nnpm run dev   # open http://localhost:5173",
        size=16,
    )
    add_text(s, Inches(0.6), Inches(4.2), Inches(12), Inches(0.4), "What you'll see:", size=15, bold=True, color=CHARCOAL)
    add_bullets(
        s,
        Inches(0.6),
        Inches(4.65),
        Inches(12),
        Inches(2.3),
        [
            "A src/ folder — your React components live here",
            "package.json — lists the libraries your app needs",
            "index.html — the one HTML page React renders into",
            "A blank Vite page in the browser — you're ready to build",
        ],
        size=16,
    )
    return s


def step_tailwind(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s, WHITE)
    add_header_band(s, "Step 2", "Style with Tailwind")
    add_text(s, Inches(0.6), Inches(1.35), Inches(12), Inches(0.4), "Tailwind = write styles as class names. No separate CSS files.", size=15, color=BROWN)
    add_code(
        s,
        Inches(0.6),
        Inches(1.9),
        Inches(12),
        Inches(1.6),
        "npm install -D tailwindcss @tailwindcss/vite\n# add @import \"tailwindcss\"; to src/index.css\n# register the plugin in vite.config.js",
        size=15,
    )
    add_text(s, Inches(0.6), Inches(3.85), Inches(12), Inches(0.4), "Then styling looks like this:", size=15, bold=True, color=CHARCOAL)
    add_code(
        s,
        Inches(0.6),
        Inches(4.3),
        Inches(12),
        Inches(2.3),
        "<button className=\"bg-amber-300 rounded-full px-4 py-2 font-bold\">\n  Sign in\n</button>",
        size=16,
    )
    return s


def step_firebase(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s, WHITE)
    add_header_band(s, "Step 3", "Create a Firebase project")
    add_bullets(
        s,
        Inches(0.6),
        Inches(1.4),
        Inches(12),
        Inches(3),
        [
            "Go to console.firebase.google.com → Add project",
            "Enable Authentication → Email/Password",
            "Create a Firestore database in test mode (for now)",
            "Register a Web app and copy the config values",
        ],
        size=17,
    )
    add_text(s, Inches(0.6), Inches(4.4), Inches(12), Inches(0.4), "Paste config into .env.local (never commit this file):", size=14, bold=True, color=CHARCOAL)
    add_code(
        s,
        Inches(0.6),
        Inches(4.85),
        Inches(12),
        Inches(2),
        "VITE_FIREBASE_API_KEY=...\nVITE_FIREBASE_AUTH_DOMAIN=huddle.firebaseapp.com\nVITE_FIREBASE_PROJECT_ID=huddle\nVITE_FIREBASE_APP_ID=...",
        size=14,
    )
    return s


def step_auth(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s, WHITE)
    add_header_band(s, "Step 4", "Add sign-up and sign-in")
    add_text(s, Inches(0.6), Inches(1.35), Inches(12), Inches(0.5), "Install the Firebase SDK, then let users create accounts.", size=15, color=BROWN)
    add_code(
        s,
        Inches(0.6),
        Inches(1.9),
        Inches(12),
        Inches(1.2),
        "npm install firebase react-router-dom",
        size=16,
    )
    add_code(
        s,
        Inches(0.6),
        Inches(3.25),
        Inches(12),
        Inches(3.5),
        "// src/services/auth.ts\nimport { createUserWithEmailAndPassword } from 'firebase/auth';\nimport { auth, db } from './firebase';\nimport { doc, setDoc } from 'firebase/firestore';\n\nexport async function signUp(name, email, password) {\n  const { user } = await createUserWithEmailAndPassword(auth, email, password);\n  await setDoc(doc(db, 'users', user.uid), { email, displayName: name });\n  return user;\n}",
        size=13,
    )
    return s


def step_firestore_model(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s, WHITE)
    add_header_band(s, "Step 5", "Design your Firestore data")
    add_text(s, Inches(0.6), Inches(1.35), Inches(12), Inches(0.5), "Four top-level collections. Each document is JSON.", size=15, color=BROWN)
    rows = [
        ("users/{uid}", "email, displayName, joinedSpotIds[], preferences", AMBER_DEEP),
        ("spots/{id}", "name, noiseLevel, outlets, lat, lng, ratingSum, ratingCount", AMBER_MID),
        ("groups/{id}", "name, course, ownerId, memberIds[], meetingTime", AMBER_DEEP),
        ("ratings/{id}", "spotId, userId, overallRating, comments, createdAt", AMBER_MID),
    ]
    for i, (name, fields, color) in enumerate(rows):
        top = Inches(2 + i * 1.1)
        pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), top, Inches(3.2), Inches(0.8))
        pill.fill.solid()
        pill.fill.fore_color.rgb = color
        pill.line.fill.background()
        pill.shadow.inherit = False
        tf = pill.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = name
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = CHARCOAL
        r.font.name = "Menlo"
        add_text(s, Inches(4), top + Inches(0.15), Inches(9), Inches(0.7), fields, size=15, color=CHARCOAL)
    return s


def step_discovery(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s, WHITE)
    add_header_band(s, "Step 6", "Show spots live from Firestore")
    add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.4), "onSnapshot streams changes in real time — no refresh needed.", size=15, color=BROWN)
    add_code(
        s,
        Inches(0.6),
        Inches(1.85),
        Inches(12),
        Inches(4.5),
        "import { collection, onSnapshot, query, orderBy } from 'firebase/firestore';\nimport { db } from './firebase';\n\nexport function subscribeSpots(cb) {\n  const q = query(collection(db, 'spots'), orderBy('name'));\n  return onSnapshot(q, (snap) =>\n    cb(snap.docs.map((d) => ({ id: d.id, ...d.data() })))\n  );\n}",
        size=14,
    )
    add_text(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4), "Tip: only subscribe after the user signs in, or Firestore rules will reject the read.", size=13, bold=True, color=BROWN)
    return s


def step_filters(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s, WHITE)
    add_header_band(s, "Step 7", "Filter chips in one array")
    add_text(s, Inches(0.6), Inches(1.35), Inches(12), Inches(0.5), "Declare filters as data, render them in a loop, AND them together.", size=15, color=BROWN)
    add_code(
        s,
        Inches(0.6),
        Inches(1.95),
        Inches(12),
        Inches(4.5),
        "const FILTERS = [\n  { key: 'quiet',   label: 'Quiet',    match: (s) => s.noiseLevel === 'Silent' },\n  { key: 'outlets', label: 'Outlets',  match: (s) => s.outlets === true },\n  { key: 'open',    label: 'Open now', match: (s) => s.open === true },\n];\n\nconst visible = spots.filter((s) =>\n  FILTERS.every((f) => !active[f.key] || f.match(s))\n);",
        size=14,
    )
    return s


def step_map(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s, WHITE)
    add_header_band(s, "Step 8", "Map with Leaflet + OpenStreetMap")
    add_code(
        s,
        Inches(0.6),
        Inches(1.35),
        Inches(12),
        Inches(1.1),
        "npm install leaflet react-leaflet",
        size=16,
    )
    add_code(
        s,
        Inches(0.6),
        Inches(2.65),
        Inches(12),
        Inches(3.8),
        "import { MapContainer, TileLayer, Marker, Popup } from 'react-leaflet';\nimport 'leaflet/dist/leaflet.css';\n\n<MapContainer center={[37.87, -122.26]} zoom={16} style={{ height: 500 }}>\n  <TileLayer url=\"https://{s}.tile.openstreetmap.org/{z}/{x}/{y}.png\" />\n  {spots.map((s) => (\n    <Marker key={s.id} position={[s.lat, s.lng]}>\n      <Popup>{s.name}</Popup>\n    </Marker>\n  ))}\n</MapContainer>",
        size=13,
    )
    return s


def step_rating(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s, WHITE)
    add_header_band(s, "Step 9", "Ratings that aggregate safely")
    add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.4), "Use a transaction so two users rating at once don't overwrite each other.", size=15, color=BROWN)
    add_code(
        s,
        Inches(0.6),
        Inches(1.9),
        Inches(12),
        Inches(4.6),
        "await runTransaction(db, async (tx) => {\n  const ref = doc(db, 'spots', spotId);\n  const snap = await tx.get(ref);\n  const data = snap.data();\n  tx.update(ref, {\n    ratingSum:   (data.ratingSum   ?? 0) + score,\n    ratingCount: (data.ratingCount ?? 0) + 1,\n  });\n  tx.set(doc(collection(db, 'ratings')), { spotId, userId, score, createdAt: serverTimestamp() });\n});",
        size=13,
    )
    return s


def step_prefs(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s, WHITE)
    add_header_band(s, "Step 10", "Let users save preferences")
    add_code(
        s,
        Inches(0.6),
        Inches(1.35),
        Inches(12),
        Inches(2.8),
        "// services/users.ts\nimport { doc, setDoc } from 'firebase/firestore';\n\nexport function updatePreferences(uid, preferences) {\n  return setDoc(doc(db, 'users', uid), { preferences }, { merge: true });\n}",
        size=14,
    )
    add_text(s, Inches(0.6), Inches(4.35), Inches(12), Inches(0.4), "merge: true means \"update these fields, keep the rest.\"", size=15, bold=True, color=CHARCOAL)
    add_text(s, Inches(0.6), Inches(4.85), Inches(12), Inches(0.4), "It also creates the document if it doesn't exist — no \"No document to update\" errors.", size=14, color=BROWN)
    return s


def step_recs(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s, WHITE)
    add_header_band(s, "Step 11", "Recommend spots by preference match")
    add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.4), "Score each spot by how many preferences it matches. Highest scores win.", size=15, color=BROWN)
    add_code(
        s,
        Inches(0.6),
        Inches(1.9),
        Inches(12),
        Inches(4.3),
        "export function scoreSpot(spot, prefs) {\n  let score = 0;\n  for (const axis of ['noiseLevel', 'outlets', 'lighting', 'crowded']) {\n    if (prefs[axis] && prefs[axis] === spot[axis]) score += 1;\n  }\n  return score;\n}\n\nconst recommended = spots\n  .map((s) => ({ s, score: scoreSpot(s, prefs) }))\n  .sort((a, b) => b.score - a.score)\n  .slice(0, 3)\n  .map(({ s }) => s);",
        size=13,
    )
    return s


def step_groups(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s, WHITE)
    add_header_band(s, "Step 12", "Study groups — create, join, leave")
    add_bullets(
        s,
        Inches(0.6),
        Inches(1.4),
        Inches(12),
        Inches(2.5),
        [
            "A group is a Firestore doc with ownerId and memberIds[]",
            "Joining = arrayUnion(userId); leaving = arrayRemove(userId)",
            "Firestore rules restrict writes to ['memberIds', 'members'] for non-owners",
        ],
        size=16,
    )
    add_code(
        s,
        Inches(0.6),
        Inches(4.2),
        Inches(12),
        Inches(2.5),
        "await updateDoc(doc(db, 'groups', groupId), {\n  memberIds: arrayUnion(uid),\n  members:   arrayUnion(displayName),\n});",
        size=15,
    )
    return s


def step_security(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s, WHITE)
    add_header_band(s, "Step 13", "Lock it down with Firestore rules")
    add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.4), "Rules live on Firebase, not in your code. They're the real security boundary.", size=14, color=BROWN)
    add_code(
        s,
        Inches(0.6),
        Inches(1.9),
        Inches(12),
        Inches(4.7),
        "match /users/{userId} {\n  allow read: if request.auth != null;\n  allow write: if request.auth.uid == userId;\n}\n\nmatch /spots/{spotId} {\n  allow read: if request.auth != null;\n  allow update: if request.resource.data\n    .diff(resource.data).affectedKeys()\n    .hasOnly(['ratingSum', 'ratingCount']);\n}",
        size=13,
    )
    return s


def step_deploy(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s, WHITE)
    add_header_band(s, "Step 14", "Ship it with Vercel")
    add_bullets(
        s,
        Inches(0.6),
        Inches(1.4),
        Inches(12),
        Inches(2),
        [
            "Push your code to GitHub",
            "Import the repo at vercel.com → 'Add New Project'",
            "Paste your VITE_FIREBASE_* values into Project → Settings → Environment Variables",
        ],
        size=17,
    )
    add_code(
        s,
        Inches(0.6),
        Inches(4),
        Inches(12),
        Inches(2),
        "npm install -g vercel\nvercel           # preview deploy\nvercel --prod    # live URL",
        size=16,
    )
    add_text(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.4), "You now have a public URL you can share with classmates.", size=14, bold=True, color=CHARCOAL)
    return s


def common_pitfalls(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s, WHITE)
    add_header_band(s, "Gotchas", "Real mistakes we hit while building this")
    bullets = [
        "Firestore IDs are STRINGS — don't parseInt() them",
        "Only subscribe to Firestore after the user signs in, or rules reject the listener and it never retries",
        "Use setDoc(..., {merge: true}) for array updates — tolerates missing docs from signup races",
        "Leaflet markers break without the icon assets — import them explicitly",
        "Don't orderBy a field that requires a composite index unless you've deployed the index",
    ]
    add_bullets(s, Inches(0.6), Inches(1.5), Inches(12), Inches(5), bullets, size=17)
    return s


def test_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s, WHITE)
    add_header_band(s, "Bonus", "Test the whole app with Playwright")
    add_code(
        s,
        Inches(0.6),
        Inches(1.35),
        Inches(12),
        Inches(1.1),
        "npm install -D @playwright/test && npx playwright install chromium",
        size=14,
    )
    add_code(
        s,
        Inches(0.6),
        Inches(2.65),
        Inches(12),
        Inches(3.5),
        "test('filter chips narrow the list', async ({ page }) => {\n  await signUpAndLand(page);\n  await page.getByRole('button', { name: 'Quiet' }).click();\n  await expect(page.getByText('Doe Library')).toBeVisible();\n  await expect(page.getByText('Cafe Strada')).toHaveCount(0);\n});",
        size=14,
    )
    return s


def next_steps(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s, AMBER_SOFT)
    add_header_band(s, "Next", "Ideas to extend Huddle")
    bullets = [
        "Upload a photo when rating a spot (Firebase Storage)",
        "Push notifications when a friend joins your group",
        "Show live occupancy counts on the map markers",
        "Swap static lat/lng for geocoded addresses",
        "Add AI spot recommendations using the rating history",
    ]
    add_bullets(s, Inches(0.6), Inches(1.5), Inches(12), Inches(4.5), bullets, size=18)
    add_text(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.6), "Each of these is its own small project — pick one and go!", size=16, bold=True, color=CHARCOAL)
    return s


def thanks_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s, NAVY)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-2), Inches(7), Inches(7))
    c.fill.solid()
    c.fill.fore_color.rgb = AMBER_DEEP
    c.line.fill.background()
    c.shadow.inherit = False
    add_text(s, Inches(0.8), Inches(2.6), Inches(9), Inches(1.5), "You built Huddle.", size=60, bold=True, color=AMBER_MID, font="Georgia")
    add_text(s, Inches(0.8), Inches(4), Inches(9), Inches(0.8), "Now ship something of your own.", size=24, color=WHITE)
    add_pill(s, Inches(0.8), Inches(5), Inches(3.2), Inches(0.55), "Questions?", fill=AMBER_DEEP, text_color=CHARCOAL, size=14)
    return s


# -------------------- Build --------------------


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = [
        title_slide,
        overview_slide,
        stack_slide,
        prereqs_slide,
        step_scaffold,
        step_tailwind,
        step_firebase,
        step_auth,
        step_firestore_model,
        step_discovery,
        step_filters,
        step_map,
        step_rating,
        step_prefs,
        step_recs,
        step_groups,
        step_security,
        step_deploy,
        common_pitfalls,
        test_slide,
        next_steps,
        thanks_slide,
    ]
    for fn in slides:
        fn(prs)

    total = len(prs.slides)
    # add footers except title/thanks
    for i, slide in enumerate(prs.slides):
        if i == 0 or i == total - 1:
            continue
        add_footer(slide, i + 1, total)

    out = "huddle-walkthrough.pptx"
    prs.save(out)
    print(f"Wrote {out} ({total} slides)")


if __name__ == "__main__":
    main()
